#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Cloud Alpacas — Final Presentation (editable DRAFT .pptx)
21 slides, 16:9 (13.333 x 7.5 in ≈ 1920x1080).

Everything is a native PowerPoint shape / text box / line — NOT a picture of a slide.
Screenshots / videos / live screens / unknown data are left as clearly-labelled
placeholder shapes for the PPT owner to replace.

Source of truth: docs 00_STORY / 04_DEMO / 07_PROPOSAL / 08_PROJECT_BRIEF /
PPT_WIREFRAME/00_WIREFRAME_GUIDE.md  (21-slide structure, Demo order, 표현 방식).
No invented facts / numbers / KPIs.
"""
import pathlib
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

OUT = pathlib.Path(__file__).parent
PPTX = OUT / "Cloud_Alpacas_Final_Presentation_Draft.pptx"

# ------------------------------------------------------------------ brand tokens
NAVY      = "0B2A47"
NAVY_DEEP = "07203A"
INK       = "1E2732"
MUTED     = "6B7787"
ORANGE    = "E77C25"
ORANGE_D  = "B85E12"
OFFWHITE  = "FAF6F0"
PAPER     = "F5F1EA"
WHITE     = "FFFFFF"
LINE      = "DAD3C7"      # warm hairline
LINE_COOL = "D8DEE7"
BLUE      = "2F6FB0"      # B2C accent
BLUE_BG   = "EAF1F8"
GREEN     = "1E7F58"      # B2B accent
GREEN_BG  = "E7F2EC"
TEAL      = "0E6E77"      # chapter III

FONT = "Pretendard"

# ------------------------------------------------------------------ geometry (in)
SW, SH = 13.333, 7.5
ML = 0.70
MR = 0.70
CW = SW - ML - MR
KICK_Y = 0.42
TITLE_Y = 0.72
RULE_Y = 1.46
BODY_Y = 1.78
BODY_B = 6.72
FOOT_Y = 6.98

prs = Presentation()
prs.slide_width  = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

# ================================================================= primitives
def slide(bg=OFFWHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = RGBColor.from_string(bg)
    r.line.fill.background()
    r.shadow.inherit = False
    return s

def _set_typeface(run, name):
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", name)

def _fmt(run, size, bold, color, name=FONT, italic=False, spacing=None):
    f = run.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.name = name
    f.color.rgb = RGBColor.from_string(color)
    _set_typeface(run, name)
    if spacing is not None:
        run._r.get_or_add_rPr().set("spc", str(int(spacing * 100)))

def text(s, x, y, w, h, content, size=16, bold=False, color=INK, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, line_spacing=1.12, name=FONT, spacing=None):
    """content: str (\\n -> paragraphs) or list of paragraphs, each str or list of (txt,bold,color)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    paras = content if isinstance(content, list) else str(content).split("\n")
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        runs = para if isinstance(para, list) else [(para, bold, color)]
        for (t, b, c) in [(r if isinstance(r, tuple) else (r, bold, color)) for r in runs]:
            run = p.add_run(); run.text = t
            _fmt(run, size, b, c, name, spacing=spacing)
    return tb

def rect(s, x, y, w, h, fill=None, line=None, lw=1.0, dash=False, round_=0.0, shadow=False):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    sh = s.shapes.add_shape(shp_type, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.shadow.inherit = False
    if round_:
        try: sh.adjustments[0] = round_
        except Exception: pass
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = RGBColor.from_string(line)
        sh.line.width = Pt(lw)
        if dash:
            ln = sh.line._get_or_add_ln()
            pd = ln.find(qn("a:prstDash"))
            if pd is None:
                pd = ln.makeelement(qn("a:prstDash"), {}); ln.append(pd)
            pd.set("val", "dash")
    if shadow:
        el = sh._element.spPr
        sp = el.makeelement(qn("a:effectLst"), {})
        o = el.makeelement(qn("a:outerShdw"),
                           {"blurRad": "90000", "dist": "40000", "dir": "5400000", "rotWithShape": "0"})
        c = el.makeelement(qn("a:srgbClr"), {"val": "0B2A47"})
        a = el.makeelement(qn("a:alpha"), {"val": "16000"})
        c.append(a); o.append(c); sp.append(o); el.append(sp)
    return sh

def hline(s, x, y, w, color=LINE, weight=1.0):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y), Inches(x + w), Inches(y))
    c.line.color.rgb = RGBColor.from_string(color); c.line.width = Pt(weight)
    c.shadow.inherit = False
    return c

# ================================================================= components
def kicker(s, txt, color=ORANGE):
    rect(s, ML, KICK_Y + 0.02, 0.14, 0.14, fill=color)
    text(s, ML + 0.26, KICK_Y - 0.06, 9.0, 0.32, txt.upper(), size=11.5, bold=True,
         color=color, spacing=2.2)

def title(s, txt, size=33, color=NAVY, rule=True):
    nlines = str(txt).count("\n") + 1
    text(s, ML, TITLE_Y, CW, 0.6 + 0.62 * nlines, txt, size=size, bold=True, color=color,
         line_spacing=1.14)
    ry = RULE_Y + 0.62 * (nlines - 1)
    if rule:
        hline(s, ML, ry, CW, LINE, 1.0)
    return ry

def footer(s, n, total=21, dark=False):
    lc = "7E92A6" if dark else MUTED
    text(s, ML, FOOT_Y, 6.0, 0.3, "CLOUD ALPACAS   ·   Final Presentation (Draft)",
         size=8.5, bold=False, color=lc, spacing=1.0)
    text(s, SW - MR - 1.6, FOOT_Y, 1.6, 0.3, f"{n:02d} / {total}", size=9.5, bold=True,
         color=lc, align=PP_ALIGN.RIGHT)

def badge(s, txt, kind, x=None, y=KICK_Y - 0.06):
    palette = {
        "PPT":         (BLUE, WHITE),
        "DEMO VIDEO":  ("5A3FA0", WHITE),
        "LIVE":        ("C0322E", WHITE),
        "PPT + 5s VIDEO": (BLUE, WHITE),
        "FORMAT TBD":  ("8A6A18", WHITE),
        "TRANSITION":  (TEAL, WHITE),
    }
    fill, tc = palette.get(kind, (MUTED, WHITE))
    w = 0.30 + 0.105 * len(txt)
    if x is None:
        x = SW - MR - w
    b = rect(s, x, y, w, 0.34, fill=fill, round_=0.5)
    text(s, x, y - 0.02, w, 0.38, txt, size=10.5, bold=True, color=tc,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=1.0)
    return x

def persona(s, x, y, name, role):
    w = 0.42 + 0.145 * (len(name) + len(role))
    rect(s, x, y, w, 0.48, fill=WHITE, line=LINE, lw=1.0, round_=0.5)
    rect(s, x + 0.09, y + 0.09, 0.30, 0.30, fill=PAPER, line=LINE, lw=0.75, round_=0.5)
    text(s, x + 0.09, y + 0.06, 0.30, 0.36, "◍", size=12, color=MUTED,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + 0.50, y, w - 0.55, 0.48,
         [[(name + "  ", True, INK), (role, False, MUTED)]],
         size=12, anchor=MSO_ANCHOR.MIDDLE)

def chips(s, x, y, items, color=INK, bg=WHITE, line=LINE, maxw=None):
    cx = x
    for it in items:
        w = 0.34 + 0.115 * len(it)
        if maxw and cx + w > x + maxw:
            cx = x; y += 0.52
        rect(s, cx, y, w, 0.40, fill=bg, line=line, lw=1.0, round_=0.5)
        text(s, cx, y - 0.01, w, 0.42, it, size=11.5, bold=True, color=color,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += w + 0.14
    return y + 0.40

ARR = "›"
def flow(s, x, y, w, nodes, h=1.0, node_h=None, big=False):
    """nodes: list of (label, sub, key)  key in {None,'b2c','b2b','hero','future','ghost'}"""
    node_h = node_h or h
    n = len(nodes)
    gap = 0.40
    nw = (w - gap * (n - 1)) / n
    for i, nd in enumerate(nodes):
        label, sub, key = (nd + (None,) * 3)[:3]
        nx = x + i * (nw + gap)
        fill, ln, tc, dash = WHITE, LINE_COOL, INK, False
        if key == "b2c":   fill, ln, tc = BLUE_BG, "AECBE6", "1C4E7E"
        elif key == "b2b": fill, ln, tc = GREEN_BG, "AAD4BF", "155E41"
        elif key == "hero": fill, ln, tc = NAVY, NAVY, WHITE
        elif key == "future": fill, ln, tc, dash = None, "4C6C89", "94A7BC", True
        elif key == "ghost": fill, ln, tc = PAPER, LINE, MUTED
        rect(s, nx, y, nw, node_h, fill=fill, line=ln, lw=1.25, dash=dash, round_=0.08,
             shadow=(key in (None,) and not big))
        if sub:
            text(s, nx + 0.06, y + 0.06, nw - 0.12, node_h - 0.44, label,
                 size=(15 if big else 12.5), bold=True, color=tc, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
            text(s, nx + 0.06, y + node_h - 0.40, nw - 0.12, 0.34, sub, size=9.5,
                 color=(tc if key in ("hero",) else MUTED), align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        else:
            text(s, nx + 0.06, y, nw - 0.12, node_h, label, size=(15 if big else 12.5),
                 bold=True, color=tc, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.05)
        if i < n - 1:
            text(s, nx + nw, y, gap, node_h, ARR, size=18, bold=True, color=MUTED,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def ph(s, x, y, w, h, label, kind="screenshot", note=None):
    """placeholder shape for the PPT owner to replace."""
    tag = {"screenshot": "SCREENSHOT", "video": "DEMO VIDEO — 재생 영역",
           "live": "LIVE — 실제 화면", "tbd": "FORMAT TBD — 표현 방식 미정",
           "data": "DATA PLACEHOLDER", "diagram": "DIAGRAM",
           "qr": "QR PLACEHOLDER", "asset": "ASSET"}.get(kind, "PLACEHOLDER")
    small = (w < 3.0 or h < 2.0)
    if kind == "video":
        rect(s, x, y, w, h, fill="1B1D22", line="3A3D44", lw=1.25, dash=True, round_=0.03)
        cx, cy = x + w / 2, y + h * 0.33
        rect(s, cx - 0.36, cy - 0.36, 0.72, 0.72, fill=None, line="C6CAD2", lw=1.75, round_=0.5)
        text(s, cx - 0.24, cy - 0.36, 0.68, 0.72, "▶", size=19, color="C6CAD2",
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, x + 0.2, y + h * 0.52, w - 0.4, 0.30, "DEMO VIDEO — 재생 영역", size=10,
             bold=True, color="AEB3BD", align=PP_ALIGN.CENTER, spacing=1.6)
        text(s, x + 0.4, y + h * 0.62, w - 0.8, 0.9, label, size=12.5, bold=True, color="E7E9EE",
             align=PP_ALIGN.CENTER, line_spacing=1.3)
        return
    elif kind == "live":
        rect(s, x, y, w, h, fill="FBEEEC", line="E0A6A3", lw=1.5, dash=True, round_=0.03)
        rect(s, x + 0.22, y + 0.22, 0.98, 0.34, fill="C0322E", round_=0.5)
        text(s, x + 0.22, y + 0.20, 0.98, 0.38, "● LIVE", size=10.5, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        tcol, tagcol = "8A3B37", "B23A35"
    elif kind == "tbd":
        rect(s, x, y, w, h, fill="FBF3E1", line="E0C57A", lw=1.5, dash=True, round_=0.03)
        tcol, tagcol = "8A6A18", "997414"
    elif kind == "qr":
        rect(s, x, y, w, h, fill=WHITE, line=INK, lw=1.5, round_=0.03)
        text(s, x, y, w, h, "QR\nPLACEHOLDER", size=13, bold=True, color=MUTED,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)
        return
    else:
        fillc = WHITE if kind in ("diagram",) else "FBFAF8"
        rect(s, x, y, w, h, fill=fillc, line=LINE, lw=1.25, dash=(kind != "diagram"), round_=0.03)
        tcol, tagcol = "6B7686", "8A94A3"
    if small:
        text(s, x + 0.1, y, w - 0.2, h, label, size=11, bold=True, color=tcol,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)
        return
    text(s, x + 0.2, y + h / 2 - 0.62, w - 0.4, 0.34, tag, size=10, bold=True,
         color=tagcol, align=PP_ALIGN.CENTER, spacing=1.6)
    text(s, x + 0.35, y + h / 2 - 0.18, w - 0.7, 0.9, label, size=13, bold=True,
         color=tcol, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.25)
    if note:
        text(s, x + 0.35, y + h - 0.62, w - 0.7, 0.5, note, size=9.5, color=tagcol,
             align=PP_ALIGN.CENTER, line_spacing=1.3)

def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt

def rail(s, x, features, ai_role, value, y=BODY_Y + 0.15, w=3.5):
    text(s, x, y, w, 0.3, "핵심 기능", size=10, bold=True, color=MUTED, spacing=1.6)
    ny = chips(s, x, y + 0.34, features, color=INK, maxw=w)
    hline(s, x, ny + 0.22, w, LINE, 1.0)
    text(s, x, ny + 0.42, w, 0.3, [[("AI 역할   ", True, NAVY), (ai_role, False, INK)]], size=12)
    text(s, x, ny + 0.80, w, 0.3, [[("Value   ", True, NAVY), (value, False, INK)]], size=12)

# ================================================================= slides
def s01():
    s = slide(NAVY_DEEP)
    rect(s, 0, 0, 0.16, SH, fill=ORANGE)
    text(s, ML + 0.2, 0.85, 10, 0.4, "CLOUD ALPACAS   ·   SALESFORCE CUSTOMER 360",
         size=12, bold=True, color=ORANGE, spacing=2.4)
    text(s, ML + 0.2, 2.35, 11.4, 2.4,
         "외부 환경에 흔들리지 않는\n지속 가능한 매출 엔진",
         size=52, bold=True, color=WHITE, line_spacing=1.12)
    text(s, ML + 0.2, 4.95, 11, 0.5, "Fan 360   →   Fan Insight   →   Partner Matching   →   Sponsorship Sales",
         size=15, color="B9C6D6")
    hline(s, ML + 0.2, 6.35, 5.6, "3C566E", 1.0)
    text(s, ML + 0.2, 6.55, 11, 0.4,
         "Cellsforce  ·  Fan Relationship Management Team          |          Final Presentation (Draft)",
         size=11, color="8FA0B4")
    ph(s, SW - 3.75, 2.35, 2.95, 2.7, "[ ALPACA MASCOT\n+ LOGO LOCKUP ]", "asset")
    notes(s, "표지. 한 문장(타이틀) + 서브 1줄 + 팀/날짜. Cloud Alpacas 로고·마스코트 자산을 우측 박스에 배치. "
             "발표자: 오늘 이야기는 '팬을 이해해 매출로 바꾸는 시스템'입니다.")

def s02():
    s = slide()
    kicker(s, "Chapter I · Overview", NAVY)
    title(s, "팬은 늘어나는데,\n왜 구단의 매출은 함께 성장하지 않을까?", size=30)
    ph(s, ML, 2.55, CW / 2 - 0.25, 3.1, "팬 수 — 우상향 (Fan Base ↑)", "data")
    ph(s, ML + CW / 2 + 0.25, 2.55, CW / 2 - 0.25, 3.1, "구단 매출 — 정체 · 적자 (Revenue ↔ / ↓)", "data")
    text(s, ML, 5.95, CW, 0.5, "팬 성장이 곧 구단의 지속 가능성은 아니다.", size=15, bold=True, color=ORANGE_D)
    notes(s, "하나의 큰 질문으로 시작. 두 그래프의 '엇갈림'이 핵심 — 실제 수치가 없으면 방향(↑ / ↔·↓)만. "
             "임의 성장률·매출 수치 만들지 않는다.")
    footer(s, 2)

def s03():
    s = slide()
    kicker(s, "Chapter I · Overview", NAVY)
    title(s, "Salesforce 도입 전, 세 가지 문제")
    data = [
        ("01", "팬 데이터가 흩어져 있다", "팬의 전체 여정을 하나로 볼 수 없다."),
        ("02", "데이터는 많지만 ACTION이 없다", "팬을 이해해도 다음 행동으로 이어지지 않는다."),
        ("03", "팬덤의 가치를 기업의 기회로 연결할 수 없다", "어떤 기업이 우리 팬덤과 맞는지 판단할 근거가 없다."),
    ]
    cw = (CW - 0.6) / 3
    for i, (num, head, sub) in enumerate(data):
        x = ML + i * (cw + 0.3)
        rect(s, x, 2.15, cw, 3.9, fill=WHITE, line=LINE, lw=1.0, round_=0.05, shadow=True)
        rect(s, x, 2.15, cw, 0.10, fill=ORANGE)
        text(s, x + 0.32, 2.5, cw - 0.6, 0.9, num, size=44, bold=True, color=NAVY)
        text(s, x + 0.32, 3.5, cw - 0.6, 1.3, head, size=17.5, bold=True, color=INK, line_spacing=1.2)
        hline(s, x + 0.32, 4.95, cw - 0.64, LINE, 1.0)
        text(s, x + 0.32, 5.12, cw - 0.6, 0.8, sub, size=11.5, color=MUTED, line_spacing=1.4)
    notes(s, "Pain Point는 정확히 3개. 큰 숫자 + 강한 headline + 짧은 subline. 긴 설명은 발표자가 말한다.")
    footer(s, 3)

def s04():
    s = slide()
    kicker(s, "Chapter I · Overview", NAVY)
    title(s, "반복은 자동화하고,\n사람은 판단과 실행에 집중한다", size=30)
    flow(s, ML, 3.05, CW,
         [("Fan Activity", "데이터 발생", None), ("Flow", "반복 자동화", None),
          ("Insight / Recommendation", "다음 행동 제안", None), ("Slack", "담당자 전달", None),
          ("Manager Action", "사람이 판단·실행", "hero")],
         node_h=1.35)
    text(s, ML, 4.85, CW, 0.4,
         "각 단계를  Salesforce · Customer 360 · Flow · Agentforce · Slack  이 자동으로 잇는다",
         size=12, color=MUTED, align=PP_ALIGN.CENTER)
    rect(s, ML, 5.55, CW, 0.9, fill=NAVY, round_=0.1)
    text(s, ML, 5.5, CW, 1.0,
         [[("DATA", True, "FFD9B4"), ("   →   ", False, "8FA0B4"), ("INSIGHT", True, "FFD9B4"),
           ("   →   ", False, "8FA0B4"), ("ACTION", True, "FFD9B4"), ("   →   ", False, "8FA0B4"),
           ("REVENUE", True, WHITE)]],
         size=18, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    notes(s, "자동화를 '기술 기능'이 아니라 '업무 방식의 변화'로 설명. 발표자: 사람이 모든 단계를 손으로 잇던 일을 "
             "플랫폼이 잇고, 사람은 판단·실행만 한다.")
    footer(s, 4)

def s05():
    s = slide()
    kicker(s, "Chapter I · Overview", NAVY)
    title(s, "B2C에서 시작해 B2B로 이어지는 하나의 프로젝트")
    flow(s, ML, 2.5, CW,
         [("B2C", "Fan Relationship · Fan Experience", "b2c"),
          ("Fan Insight", "연결점 · Bridge", "b2c"),
          ("B2B", "Partner Matching · Sponsorship Sales", "b2b")],
         node_h=1.5)
    text(s, ML, 4.5, CW, 0.3, "Feature Owner", size=10, bold=True, color=MUTED, spacing=1.8)
    chips(s, ML, 4.85,
          ["사라 · Fan 360 / Insight", "승우 · Product · Quote · Campaign", "은영 · Opportunity",
           "혜준 · Lead", "아론 · Account · Contact"], color=INK, maxw=CW)
    text(s, ML, 5.7, CW, 0.4,
         "각자 자기 구간을 하나의 작은 Salesforce 프로젝트처럼 책임 — 하나의 시나리오로 연결.",
         size=12, color=MUTED)
    notes(s, "B2C와 B2B가 별개 프로젝트가 아니라 Fan Insight로 이어진 하나의 흐름임을 보여준다. 팀은 1줄로.")
    footer(s, 5)

def s06():
    s = slide()
    kicker(s, "Chapter II · Demo Scenario", ORANGE)
    title(s, "From Fan Action to Sponsorship Revenue")
    text(s, ML, RULE_Y + 0.12, CW, 0.35, "지금부터 보실 Demo의 전체 지도 — 10초 안에.",
         size=13, color=MUTED)
    flow(s, ML, 3.2, CW,
         [("Fan", "QR 참여", "b2c"), ("Fan 360", "팬 상태", "b2c"), ("Fan Insight", "기회 발견", "b2c"),
          ("Partner\nMatching", "기업 매칭", "b2b"), ("Lead", "영업 대상", "b2b"),
          ("Opportunity", "딜", "b2b"), ("AI Sales", "가속", "b2b"), ("Closed Won", "매출", "b2b")],
         node_h=1.25)
    text(s, ML, 5.35, CW, 0.4,
         "B2C(파랑) 3단계  →  Fan Insight 연결점  →  B2B(초록) 5단계",
         size=12, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    notes(s, "Demo 목차. 'Business Opportunity'라는 말은 쓰지 않는다 — Partner Matching. Fan Insight가 B2C의 "
             "마지막이자 B2B의 출발점.")
    footer(s, 6)

def s07():
    s = slide(NAVY_DEEP)
    badge(s, "LIVE", "LIVE", x=ML, y=0.5)
    text(s, ML, 1.15, 11.5, 0.9, "GAME DAY LIVE", size=42, bold=True, color=WHITE, spacing=1.5)
    text(s, ML, 2.1, 11.5, 0.5, "⚾  7회말 경기 진행 중", size=16, color="9DB0C4")
    rect(s, ML, 2.95, CW, 2.75, fill="0E2B3A", line=ORANGE, lw=1.5, round_=0.04)
    text(s, ML + 0.55, 3.35, CW - 4.2, 1.6,
         "오늘의 FAN EVENT — 문태양 선수 퀴즈",
         size=26, bold=True, color=WHITE, line_spacing=1.25)
    text(s, ML + 0.55, 4.75, CW - 4.2, 0.5, "QR을 스캔하고 퀴즈에 참여하세요", size=15, color="FFD9B4")
    ph(s, SW - MR - 2.55, 3.25, 2.15, 2.15, "", "qr")
    flow(s, ML, 6.0, CW,
         [("관객 QR 참여", "", "ghost"), ("Quiz Entry", "", "ghost"), ("Fan Activity", "", "ghost"),
          ("Salesforce", "", "ghost"), ("Fan 360", "", "ghost")], node_h=0.6)
    notes(s, "실제 경기장 전광판처럼. QR destination이 확정 안 됐으면 'QR PLACEHOLDER' 유지 — 가짜 QR을 실제처럼 "
             "표현하지 않는다. Cloud Alpacas 알파카 마스코트를 이벤트 안내 그래픽으로 활용 가능. "
             "발표자: 지금 여러분은 경기장 관객입니다. 전광판에 이벤트가 떴습니다. QR을 찍어보세요. "
             "→ 이 참여가 곧 CRM 데이터가 됩니다. (04_DEMO Scene 1 = 'PPT + 관객 참여')")
    footer(s, 7, dark=True)

def demo(n, chap_b2b, badge_kind, question, persona_nm, persona_role, ph_label, ph_kind,
         features, ai_role, value, note, ph_note=None):
    s = slide()
    kicker(s, "Chapter II · Demo — " + ("B2B" if chap_b2b else "B2C"), GREEN if chap_b2b else BLUE)
    badge(s, badge_kind, badge_kind)
    if persona_nm:
        persona(s, ML, KICK_Y + 0.42, persona_nm, persona_role)
    ty = 1.52
    nl = str(question).count("\n") + 1
    text(s, ML, ty, CW - 3.7, 0.6 + 0.55 * nl, question, size=28, bold=True, color=NAVY, line_spacing=1.14)
    ry = ty + 0.62 + 0.62 * (nl - 1)
    hline(s, ML, ry, CW, LINE, 1.0)
    py = ry + 0.30
    pw = CW - 3.9
    ph(s, ML, py, pw, BODY_B - py - 0.05, ph_label, ph_kind, note=ph_note)
    rail(s, ML + pw + 0.45, features, ai_role, value, y=py + 0.05, w=3.45)
    notes(s, note)
    footer(s, n)
    return s

def s08():
    demo(8, False, "PPT", "우리 팬은 누구인가?", "김매니저", "FRM Manager",
         "[ SCREENSHOT ]\nFan 360  →  Segment  →  Recommendation Hub 핵심 화면",
         "screenshot", ["Fan 360", "Segment", "Recommendation Hub"],
         "팬 이해 지원", "Fan Understanding",
         "발표자: 티켓·굿즈·멤버십·참여가 흩어져 있던 데이터를 한 팬의 프로필·타임라인으로. "
         "실제 Fan 360 / Fan Profile 화면 캡처. (04_DEMO Scene 2 앞부분)")

def s09():
    demo(9, False, "DEMO VIDEO", "각 팬에게 어떻게 다르게 행동할까?", "김매니저", "FRM Manager",
         "Target Fan 확인  →  AI 개인화 메시지 생성  →  Review  →  발송\n(약 80–90초)",
         "video", ["AI Personalized Message"], "Personalize", "Personalized Fan Engagement",
         "영상 재생 영역이 화면의 중심. 발표자: 팬별 특성을 담당자가 일일이 쓰기 어려웠다 — AI가 생성하고 "
         "담당자가 검토·승인 후 발송(Human-in-the-loop). (04_DEMO 데모 영상 ①)")

def s10():
    s = slide()
    kicker(s, "Chapter II · Demo — Bridge", TEAL)
    badge(s, "TRANSITION", "TRANSITION")
    title(s, "팬의 관심사가\n기업을 찾는 근거가 된다", size=29)
    y = 3.15
    hh = 2.05
    lw = CW / 2 - 0.6
    pill_w = 2.1
    rect(s, ML, y, lw, hh, fill=BLUE_BG, line="AECBE6", lw=1.25, round_=0.07)
    text(s, ML + 0.34, y + 0.26, lw - 0.68, 0.3, "B2C · 여기까지", size=11, bold=True, color=BLUE, spacing=1.4)
    text(s, ML + 0.34, y + 0.72, lw - 1.0, 1.1,
         [["Fan Data", ("  →  ", False, MUTED), "Fan 360", ("  →  ", False, MUTED), "Fan Insight"]],
         size=14, bold=True, color="1C4E7E", line_spacing=1.35)
    rx = ML + CW / 2 + 0.6
    rect(s, rx, y, lw, hh, fill=GREEN_BG, line="AAD4BF", lw=1.25, round_=0.07)
    text(s, rx + 0.34, y + 0.26, lw - 0.68, 0.3, "B2B · 여기부터", size=11, bold=True, color=GREEN, spacing=1.4)
    text(s, rx + 0.62, y + 0.72, lw - 0.9, 1.1,
         [["Partner Matching", ("  →  ", False, MUTED), "Sponsorship Sales"]],
         size=14, bold=True, color="155E41", line_spacing=1.35, align=PP_ALIGN.RIGHT)
    rect(s, SW / 2 - pill_w / 2, y + hh / 2 - 0.44, pill_w, 0.88, fill=WHITE, line=NAVY, lw=1.75,
         round_=0.5, shadow=True)
    text(s, SW / 2 - pill_w / 2, y + hh / 2 - 0.46, pill_w, 0.9, "FAN INSIGHT", size=15, bold=True,
         color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, ML, 5.7, CW, 0.5,
         "20·30대 여성 팬 증가  ·  뷰티 관심 Signal        →        Partner Matching의 근거",
         size=14, bold=True, color=ORANGE_D, align=PP_ALIGN.CENTER)
    notes(s, "Demo 전체의 turning point. 여기서 색·발표 톤이 파랑(B2C)→초록(B2B)으로 전환. Fan Insight는 B2B 안이 "
             "아니라 B2C의 마지막 단계. 발표자: 팬의 관심사가 곧 어떤 기업을 찾을지의 근거가 됩니다.")
    footer(s, 10)

def s11():
    demo(11, True, "LIVE", "팬 데이터를 어떻게\nB2B 영업 기회로 연결할까?", "김매니저 → 이매니저", "FRM → Sponsorship Sales",
         "Monthly Fan Insight Letter 확인  →  Slack Agent 분석 요청  →  적합한 Sponsorship 방향 탐색",
         "live", ["Monthly Fan Insight Letter", "Slack Agent"], "Analyze & Discover",
         "B2C Data → B2B Sales Opportunity",
         "현장 시연. 발표자: B2C에서 쌓인 팬 데이터를 B2B 담당자가 직접 다시 해석해야 했다 — 이제 월간 Fan Insight "
         "Letter와 Slack Agent가 방향을 좁혀준다. 김매니저 → 이매니저 handoff. (04_DEMO Scene 3)")

def s12():
    s = demo(12, True, "FORMAT TBD", "이 팬덤과 가장 잘 맞는 기업은?\n— 왜 이 기업인가?", "이매니저", "Sponsorship Sales Manager",
             "[ FORMAT TBD ]\nFan Insight  →  Fan Fit  →  Partner Candidate  →  d'Alba  →  Sales Opportunity",
             "tbd", ["Fan Fit", "Segment Match", "Lead Score"], "Match & Explain", "데이터 기반 파트너 발굴",
             "‘d'Alba는 팬덤과 높은 적합도를 보여 후보가 됐습니다’ — AI는 정답이 아니라 '왜 이 기업인가'를 설명. "
             "Fan Fit / Segment Match(팬덤 적합도) ≠ Lead Score(계약 가능성). 기업 데이터 = OpenDART API. "
             "※ Demo순서 이미지에 S4 열이 없어 재구성 — 포함 여부·표현 방식 팀 확정 필요.",
             ph_note="Demo순서 이미지에 S4 미표시 — PPT / VIDEO / LIVE 중 미선택")
    return s

def s13():
    demo(13, True, "PPT + 5s VIDEO", "Sponsor 후보를 어떻게\n실제 Deal로 발전시킬까?", "이매니저 / d'Alba", "Sponsorship Sales",
         "[ SCREENSHOT  +  ▶ 5s ]\nTableau Next Dashboard  →  Lead / Lead Score  →  Account (AI Enrichment)  →  d'Alba OPP",
         "screenshot", ["Tableau Next", "Lead Score", "Account AI Enrichment"],
         "Analyze / Score / Enrich", "Sales Prioritization · Productivity · Data Quality",
         "PPT 위에 약 5초 임베드 영상. 발표자: 유망 Sponsor 판단부터 Account 정보 보완까지 수작업이 많았다 — "
         "우선순위는 Lead Score, 빈 정보는 AI가 DART로 채운다. (04_DEMO Scene 4)")

def s14():
    demo(14, True, "DEMO VIDEO", "고객은 무엇을 말했는가?", "이매니저 / 김하나", "Sales / d'Alba 담당자",
         "고객 Meeting / Activity  →  기록  →  AI 분석  →  Summary / Signal\n(약 3분)",
         "video", ["Activity Intelligence"], "Understand", "Activity 자산화",
         "영상 재생 영역이 화면의 중심. 발표자: 미팅·대화가 단순 기록으로 남아 다시 읽고 해석해야 했다 — 이제 대화가 "
         "요약·Signal(긍정/위험)로 정리된다. Zoom 연동 · Interaction Intelligence. (04_DEMO Scene 6)")

def s15():
    demo(15, True, "LIVE", "그래서 무엇을 제안할까?", "이매니저 / d'Alba", "Sponsorship Sales",
         "[ LIVE — Opportunity Agent ]\n과거 유사 사례 + 현재 OPP + 팬 데이터 + 고객 Activity  →  제안 방향 / Package / Product + 근거",
         "live", ["Opportunity Agent"], "Reason", "Context 기반 Sales Decision Support",
         "현장 시연. 발표자: 담당자가 컨텍스트를 직접 모아 판단해야 했다 — Agent가 근거와 함께 제안 방향을 제시. "
         "조회·추천은 즉시, 쓰기는 담당자 승인 후. (04_DEMO Scene 7)")

def s16():
    demo(16, True, "PPT", "고객의 변화에\n어떻게 대응할까?", "이매니저 / 김하나", "Sales / d'Alba 담당자",
         "[ SCREENSHOT ]\n새 고객 Activity / 상황  →  AI 선제 분석  →  Negotiation 대응 / 수정안 + 근거  →  Closed Won",
         "screenshot", ["Proactive AI", "Negotiation Assistant"], "Act Proactively", "Proactive Selling",
         "발표자: 고객 반응 변화마다 다시 분석해야 했다 — AI가 변화를 먼저 감지해 수정안을 근거와 함께 제시. "
         "최종 결정·승인은 담당자. ⚠️ 구체 금액은 발표 전 하나로 통일 (04_DEMO 가격 검증). (04_DEMO Scene 8)")

def s17():
    demo(17, True, "FORMAT TBD", "1년 후, 관계를\n어떻게 다음 매출로 연결할까?", "이매니저 / d'Alba", "Sponsorship Sales",
         "[ FORMAT TBD ]\n‘1년 후’  →  단년 계약 종료 임박  →  Partnership Plan 확인  →  d'Alba Upsell",
         "tbd", ["Partnership Plan (논의 필요)", "Upsell (논의 필요)"], "미정",
         "Renewal / Upsell Revenue Expansion (방향)",
         "Demo순서 이미지: 표현 방식·AI 역할·기능 모두 미정. 재계약·장기 Partnership 자동화는 문서상 Future Scope — "
         "구현 확인된 부분만 시연. 발표자: '첫 계약은 매출엔진의 끝이 아니라 시작' — 단, 여기부터는 방향 제시. (04_DEMO Scene 9)",
         ph_note="AI 역할·기능·표현 방식 모두 미정 — 팀 확정 필요")

def s18():
    s = slide()
    kicker(s, "Chapter III · Global Best Practices", TEAL)
    title(s, "우리가 이렇게 설계한 이유")
    data = [
        ("STANDARD FIRST", "표준 Salesforce를 먼저 활용하고, 필요한 부분만 Custom."),
        ("ONE CUSTOMER VIEW", "Fan 360을 중심으로 B2C와 B2B가 같은 데이터를 활용."),
        ("AI + AUTOMATION", "반복은 자동화, AI는 제안, 최종 판단은 사람이 한다."),
    ]
    cw = (CW - 0.6) / 3
    for i, (head, sub) in enumerate(data):
        x = ML + i * (cw + 0.3)
        rect(s, x, 2.1, cw, 2.5, fill=WHITE, line=LINE, lw=1.0, round_=0.05, shadow=True)
        text(s, x + 0.3, 2.42, cw - 0.6, 0.5, f"0{i+1}", size=24, bold=True, color=TEAL)
        text(s, x + 0.3, 3.0, cw - 0.6, 0.7, head, size=16.5, bold=True, color=INK, spacing=0.6)
        text(s, x + 0.3, 3.62, cw - 0.6, 0.9, sub, size=11.5, color=MUTED, line_spacing=1.4)
    text(s, ML, 4.95, CW, 0.28, "전체 흐름", size=9.5, bold=True, color=MUTED, spacing=1.8)
    flow(s, ML, 5.28, CW,
         [("Fan Activity", "", None), ("Data", "", None), ("Insight", "", None),
          ("AI / Automation", "", None), ("Human Action", "", "hero")], node_h=0.7)
    text(s, ML, 6.15, CW, 0.3,
         "17 Custom Objects · 40 Flows · 1 Trigger · 46 LWC · 5 Agentforce Agents · 6 Prompt Templates  (실측)",
         size=9.5, color=MUTED)
    notes(s, "4개 기술 나열이 아니라 설계 원칙 3개. 각 원칙은 Headline + 1줄. 구현 수치는 발표자가 필요할 때만 언급. "
             "'왜 이렇게 설계했는가'를 설명하는 페이지.")
    footer(s, 18)

def s19():
    s = slide()
    kicker(s, "Chapter III · Global Best Practices", TEAL)
    title(s, "설계 원칙이 실제 업무 흐름으로 작동한다")
    flow(s, ML, 2.7, CW,
         [("Event / Fan Activity", "데이터 발생", None), ("Flow", "반복 자동화", None),
          ("Recommendation", "다음 행동 제안", None), ("Slack", "담당자 전달", None),
          ("Manager Action", "사람이 판단·실행", "hero")],
         node_h=1.25)
    text(s, ML, 4.6, CW, 0.4,
         "Salesforce Platform · Customer 360 · Flow · Agentforce · Slack  이 커넥터처럼 잇는다",
         size=12, color=MUTED, align=PP_ALIGN.CENTER)
    rect(s, ML, 5.25, CW, 1.15, fill=NAVY, round_=0.09)
    text(s, ML + 0.5, 5.25, CW - 1.0, 1.15,
         "자동화의 목적은 사람을 대체하는 것이 아니라,\n담당자가 판단하고 실행할 시간을 확보하는 것.",
         size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)
    notes(s, "기존 Automation 페이지 중심. 단계별 1줄 역할. 긴 설명 금지. 발표자: Flow가 반복을 처리하고 사람은 "
             "판단·실행에 집중한다.")
    footer(s, 19)

def s20():
    s = slide()
    kicker(s, "Chapter IV · Conclusion", NAVY)
    title(s, "우리가 만든 것")
    data = [
        ("FAN", "팬을 이해하다", ["Fan 360", "Personalization"], "Fan Lifetime Value  ↑"),
        ("INSIGHT", "기회를 발견하다", ["Fan Insight", "Partner Matching"], "Personalized Fan Experience  ↑"),
        ("REVENUE", "매출로 연결하다", ["Sponsorship Sales", "Opportunity"], "Sponsorship Revenue  ↑"),
    ]
    cw = (CW - 1.1) / 3
    for i, (kw, tag, impl, val) in enumerate(data):
        x = ML + i * (cw + 0.55)
        rect(s, x, 2.15, cw, 4.0, fill=WHITE, line=LINE, lw=1.0, round_=0.05, shadow=True)
        text(s, x, 2.55, cw, 0.8, kw, size=34, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        text(s, x, 3.35, cw, 0.4, tag, size=13, color=MUTED, align=PP_ALIGN.CENTER)
        cy = 4.0
        for j, a in enumerate(impl):
            w = 0.34 + 0.115 * len(a)
            rect(s, x + (cw - w) / 2, cy + j * 0.5, w, 0.4, fill=PAPER, line=LINE, lw=1.0, round_=0.5)
            text(s, x, cy + j * 0.5 - 0.01, cw, 0.42, a, size=11, bold=True, color=INK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        hline(s, x + 0.3, 5.55, cw - 0.6, LINE, 1.0)
        text(s, x, 5.72, cw, 0.4, val, size=13.5, bold=True, color=ORANGE_D, align=PP_ALIGN.CENTER)
        if i < 2:
            text(s, x + cw, 2.15, 0.55, 4.0, ARR, size=20, bold=True, color=MUTED,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, ML, 6.35, CW, 0.35, "FAN  →  INSIGHT  →  REVENUE — 하나의 흐름", size=12, bold=True,
         color=MUTED, align=PP_ALIGN.CENTER)
    notes(s, "What We Built + Business Value 통합. 한 페이지에서 FAN→INSIGHT→REVENUE가 한눈에. 미측정 KPI·ROI 금지.")
    footer(s, 20)

def s21():
    s = slide(NAVY_DEEP)
    rect(s, 0, 0, 0.16, SH, fill=ORANGE)
    text(s, ML + 0.2, 0.55, 6, 0.3, "NOW · 현재 구현", size=10.5, bold=True, color="9DB0C4", spacing=1.8)
    flow(s, ML + 0.2, 0.95, CW - 0.4,
         [("Fan Data", "", "b2c"), ("Insight", "", "b2c"), ("Partner Matching", "", "b2b"),
          ("Sponsorship", "", "b2b")], node_h=0.62)
    text(s, ML + 0.2, 1.9, 6, 0.3, "FUTURE SCOPE · 미구현", size=10.5, bold=True, color="D6A6A3", spacing=1.8)
    flow(s, ML + 0.2, 2.3, CW - 0.4,
         [("Real-time Data", "", "future"), ("AI Decision", "", "future"),
          ("Autonomous Action", "", "future"), ("Continuous Revenue Growth", "", "future")], node_h=0.62)
    text(s, ML + 0.2, 3.75, 8.7, 1.9,
         "팬을 이해하고, 팬덤의 가치를 발견하고,\n그 가치를 매출로 연결합니다.",
         size=33, bold=True, color=WHITE, line_spacing=1.3)
    text(s, ML + 0.2, 5.95, 8.5, 0.5,
         [[("CLOUD ALPACAS", True, ORANGE), ("      ·      Sustainable Revenue Engine", False, "B9C6D6")]],
         size=15)
    ph(s, SW - MR - 2.7, 3.9, 2.7, 1.85, "[ ALPACA + LOGO ]", "asset")
    notes(s, "발표의 마지막 슬라이드. FUTURE는 현재 구현처럼 보이지 않게(점선·흐리게). 별도 Closing/당첨자 슬라이드 "
             "없음 — 이후 발표자가 '오늘 참여해주신 분들 중…' 하며 LIVE 당첨자 발표 → Q&A로 전환.")
    footer(s, 21, dark=True)

for fn in [s01, s02, s03, s04, s05, s06, s07, s08, s09, s10, s11, s12, s13, s14, s15, s16, s17, s18, s19, s20, s21]:
    fn()

prs.save(str(PPTX))
print(f"saved {PPTX.name}  ({len(prs.slides._sldIdLst)} slides)")
