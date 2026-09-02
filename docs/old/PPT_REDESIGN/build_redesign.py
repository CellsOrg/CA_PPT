#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Cloud Alpacas — Final Presentation  ·  REDESIGN (editable .pptx)

같은 콘텐츠 · 같은 페이지 순서(21장)를 유지하고, Visual Design / Layout /
Typography / Information Hierarchy 만 처음부터 다시 설계한 버전.

Design language : Editorial + Modern Corporate + Minimal Sales Deck
 - 넓은 macro-whitespace, 한 페이지 = 하나의 핵심 메시지
 - 큰 headline typography (Pretendard weight contrast), 짧은 문장 + 큰 키워드
 - 얇은 hairline + subtle divider, 카드/박스 최소화
 - Warm off-white canvas · Deep Navy · 절제된 Orange accent
 - 필요한 경우에만 diagram, 이미지/영상/LIVE 영역은 크게 하나만

Content source of truth : docs 00_STORY / 04_DEMO / 07_PROPOSAL / 08_PROJECT_BRIEF /
PPT_WIREFRAME/00_WIREFRAME_GUIDE.md 및 기존 PPT_DRAFT/build_pptx.py.
문구·수치·페이지 순서는 그대로. 새 사실/숫자/KPI를 만들지 않는다.
"""
import pathlib
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

OUT = pathlib.Path(__file__).parent
PPTX = OUT / "Cloud_Alpacas_Presentation_Redesign.pptx"

# ------------------------------------------------------------------ brand tokens
NAVY      = "0B2A47"
NAVY_DEEP = "071B2E"
INK       = "20272F"
MUTED     = "8B8E93"
FAINT     = "B7B4AD"
ORANGE    = "E77C25"
ORANGE_D  = "B85E12"
PAPER     = "FBFAF8"     # warm off-white canvas
PAPER_2   = "F4F1EB"     # slightly deeper warm
WHITE     = "FFFFFF"
LINE      = "E5E1D9"     # warm hairline
LINE_DK   = "31465C"     # hairline on navy
DARK_VID  = "14171B"

# Pretendard weight faces (installed OTF family names)
F_LIGHT = "Pretendard Light"
F_REG   = "Pretendard"
F_MED   = "Pretendard Medium"
F_SB    = "Pretendard SemiBold"
F_BOLD  = "Pretendard"   # + bold=True

# ------------------------------------------------------------------ geometry (in)
SW, SH = 13.333, 7.5
ML = 0.92
CW = SW - 2 * ML
KICK_Y  = 0.60
HEAD_Y  = 1.42
FOOT_Y  = 7.02
TOTAL   = 21

prs = Presentation()
prs.slide_width  = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

# ================================================================= primitives
def slide(bg=PAPER):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = RGBColor.from_string(bg)
    r.line.fill.background(); r.shadow.inherit = False
    return s

def _typeface(run, name):
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", name)

def _fmt(run, size, color, name=F_REG, bold=False, italic=False, spacing=None):
    f = run.font
    f.size = Pt(size); f.bold = bold; f.italic = italic; f.name = name
    f.color.rgb = RGBColor.from_string(color)
    _typeface(run, name)
    if spacing is not None:
        run._r.get_or_add_rPr().set("spc", str(int(spacing * 100)))

def text(s, x, y, w, h, content, size=15, name=F_REG, color=INK, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, line_spacing=1.16, bold=False, spacing=None):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    paras = content if isinstance(content, list) else str(content).split("\n")
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line_spacing
        runs = para if isinstance(para, list) else [(para, name, color)]
        for item in runs:
            t, nm, col = (item if isinstance(item, tuple) else (item, name, color))
            run = p.add_run(); run.text = t
            _fmt(run, size, col, nm, bold=bold, spacing=spacing)
    return tb

def rect(s, x, y, w, h, fill=None, line=None, lw=1.0, dash=False):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.shadow.inherit = False
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = RGBColor.from_string(line); sh.line.width = Pt(lw)
        if dash:
            ln = sh.line._get_or_add_ln()
            pd = ln.find(qn("a:prstDash")) or ln.makeelement(qn("a:prstDash"), {})
            if pd.getparent() is None: ln.append(pd)
            pd.set("val", "dash")
    return sh

def oval(s, x, y, d, fill=None, line=None, lw=1.0):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    sh.shadow.inherit = False
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb = RGBColor.from_string(line); sh.line.width = Pt(lw)
    return sh

def hline(s, x, y, w, color=LINE, weight=0.75, dash=False):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y), Inches(x + w), Inches(y))
    c.line.color.rgb = RGBColor.from_string(color); c.line.width = Pt(weight)
    c.shadow.inherit = False
    if dash:
        ln = c.line._get_or_add_ln()
        pd = ln.makeelement(qn("a:prstDash"), {}); pd.set("val", "sysDot"); ln.append(pd)
    return c

def vline(s, x, y, h, color=LINE, weight=0.75):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y), Inches(x), Inches(y + h))
    c.line.color.rgb = RGBColor.from_string(color); c.line.width = Pt(weight)
    c.shadow.inherit = False
    return c

def diag(s, x1, y1, x2, y2, color=INK, weight=1.5):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = RGBColor.from_string(color); c.line.width = Pt(weight)
    c.shadow.inherit = False
    return c

def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt

# ================================================================= chrome
def kicker(s, txt, color=NAVY, dark=False):
    rect(s, ML, KICK_Y + 0.015, 0.11, 0.11, fill=color)
    text(s, ML + 0.24, KICK_Y - 0.08, 9.0, 0.32, txt.upper(), size=10.5, name=F_SB,
         color=(WHITE if dark else color), spacing=2.6)

def badge(s, txt, dark=False):
    text(s, SW - ML - 4.2, KICK_Y - 0.08, 4.2, 0.32, txt.upper(), size=10, name=F_SB,
         color=("9DB0C4" if dark else MUTED), align=PP_ALIGN.RIGHT, spacing=2.4)

def headline(s, txt, y=HEAD_Y, size=35, color=NAVY, weight=F_MED, w=CW, rule=True):
    n = str(txt).count("\n") + 1
    text(s, ML, y, w, 0.6 + 0.66 * n, txt, size=size, name=weight, color=color, line_spacing=1.18)
    ry = y + 0.30 + 0.66 * n
    if rule:
        hline(s, ML, ry, CW, LINE, 0.75)
    return ry

def footer(s, n, dark=False):
    col = "6E8093" if dark else FAINT
    text(s, ML, FOOT_Y, 5.0, 0.3, "CLOUD ALPACAS", size=8, name=F_SB, color=col, spacing=2.2)
    text(s, SW - ML - 1.6, FOOT_Y, 1.6, 0.3, f"{n:02d} / {TOTAL}", size=8.5, name=F_SB,
         color=col, align=PP_ALIGN.RIGHT, spacing=1.4)

def alpaca(s, x, y, w, h, label="ALPACA  +  LOGO", dark=False):
    rect(s, x, y, w, h, fill=None, line=(LINE_DK if dark else LINE), lw=1.0, dash=True)
    text(s, x, y, w, h, label, size=9.5, name=F_SB, color=("6E8093" if dark else FAINT),
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=2.0)

# ================================================================= diagram helpers
ARR = "→"

def flow(s, x, y, w, nodes, accent_last=True, sub=True, dark=False,
         label_size=11.5, dot_at=None):
    """nodes: list of (label, sublabel) or (label,).  Minimal: hairline + hollow dots."""
    base = "43566A" if dark else LINE
    tcol = WHITE if dark else INK
    scol = "7E92A6" if dark else MUTED
    n = len(nodes)
    hline(s, x + 0.1, y, w - 0.2, base, 0.9)
    step = (w - 0.2) / (n - 1) if n > 1 else 0
    for i, nd in enumerate(nodes):
        label = nd[0]; slab = nd[1] if (len(nd) > 1) else ""
        cx = x + 0.1 + i * step
        last = (i == n - 1)
        emph = accent_last and last
        d = 0.16
        if emph:
            oval(s, cx - d / 2, y - d / 2, d, fill=ORANGE)
        else:
            oval(s, cx - d / 2, y - d / 2, d, fill=(NAVY_DEEP if dark else PAPER),
                 line=(WHITE if dark else NAVY), lw=1.2)
        tw = step if 0 < i < n - 1 else step * 0.9
        text(s, cx - tw / 2, y - 0.66, tw, 0.5, label,
             size=label_size, name=(F_SB if emph else F_MED),
             color=(ORANGE if emph else tcol),
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM, line_spacing=1.08)
        if sub and slab:
            text(s, cx - tw / 2, y + 0.18, tw, 0.42, slab, size=9, name=F_REG,
                 color=scol, align=PP_ALIGN.CENTER, line_spacing=1.05)

def media(s, x, y, w, h, kind, label, note=None):
    """Large single media / screen area — the hero of demo pages."""
    cx = x + w / 2
    if kind == "video":
        rect(s, x, y, w, h, fill=DARK_VID, line="2C3138", lw=1.0)
        oval(s, cx - 0.34, y + h * 0.30 - 0.34, 0.68, fill=None, line="C9CDD4", lw=1.5)
        text(s, cx - 0.30, y + h * 0.30 - 0.34, 0.72, 0.68, "▶", size=17,
             color="C9CDD4", align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, x + 0.4, y + h * 0.30 + 0.52, w - 0.8, 0.3, "DEMO VIDEO", size=9.5,
             name=F_SB, color="8C929B", align=PP_ALIGN.CENTER, spacing=3.0)
        text(s, x + 0.8, y + h * 0.30 + 0.98, w - 1.6, 1.0, label, size=13, name=F_MED,
             color="E4E6EA", align=PP_ALIGN.CENTER, line_spacing=1.35)
        return
    if kind == "live":
        rect(s, x, y, w, h, fill=WHITE, line=LINE, lw=1.0)
        oval(s, x + 0.34, y + 0.36, 0.13, fill=ORANGE)
        text(s, x + 0.56, y + 0.28, 2.0, 0.3, "LIVE", size=10, name=F_SB, color=ORANGE_D,
             spacing=2.6)
        text(s, x + 0.8, y + h / 2 - 0.5, w - 1.6, 1.2, label, size=13, name=F_MED,
             color=INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.4)
        text(s, x + 0.4, y + h - 0.5, w - 0.8, 0.3, "발표 중 실제 화면으로 전환",
             size=9, name=F_REG, color=FAINT, align=PP_ALIGN.CENTER)
        return
    if kind == "tbd":
        rect(s, x, y, w, h, fill=None, line=FAINT, lw=1.0, dash=True)
        text(s, x + 0.4, y + h / 2 - 0.72, w - 0.8, 0.3, "FORMAT TBD", size=10, name=F_SB,
             color=MUTED, align=PP_ALIGN.CENTER, spacing=3.0)
        text(s, x + 0.8, y + h / 2 - 0.28, w - 1.6, 1.0, label, size=13, name=F_MED,
             color=INK, align=PP_ALIGN.CENTER, line_spacing=1.4)
        if note:
            text(s, x + 0.8, y + h - 0.6, w - 1.6, 0.4, note, size=9, name=F_REG,
                 color=FAINT, align=PP_ALIGN.CENTER)
        return
    # screenshot (default), optionally with 5s inset
    rect(s, x, y, w, h, fill=WHITE, line=LINE, lw=1.0, dash=True)
    text(s, x + 0.4, y + h / 2 - 0.62, w - 0.8, 0.3,
         ("SCREENSHOT  +  ▶ 5s" if kind == "shot5s" else "SCREENSHOT"),
         size=10, name=F_SB, color=MUTED, align=PP_ALIGN.CENTER, spacing=3.0)
    text(s, x + 0.8, y + h / 2 - 0.18, w - 1.6, 1.0, label, size=13, name=F_MED,
         color=INK, align=PP_ALIGN.CENTER, line_spacing=1.4)
    if kind == "shot5s":
        rect(s, x + w - 1.5, y + h - 0.78, 1.2, 0.5, fill=DARK_VID)
        text(s, x + w - 1.5, y + h - 0.80, 1.2, 0.5, "▶ 5s", size=9.5, name=F_SB,
             color="D8DADE", align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ================================================================= DEMO template
def demo(n, side, badge_txt, persona, question, media_kind, media_label, meta,
         note, media_note=None):
    s = slide()
    acc = NAVY if side == "B2C" else ORANGE
    kicker(s, f"Chapter II · Demo — {side}", acc)
    badge(s, badge_txt)
    text(s, ML + 0.24, KICK_Y + 0.30, 8.0, 0.3, persona, size=10.5, name=F_MED, color=MUTED)
    ry = headline(s, question, y=1.30, size=31, weight=F_MED)
    my = ry + 0.28
    mh = 6.34 - my
    media(s, ML, my, CW, mh, media_kind, media_label, note=media_note)
    text(s, ML, 6.48, CW, 0.34, meta, size=9.5, name=F_MED, color=MUTED, spacing=0.4)
    notes(s, note)
    footer(s, n)
    return s

# ================================================================= slides
def s01():
    s = slide(NAVY_DEEP)
    rect(s, 0, 0, 0.10, SH, fill=ORANGE)
    text(s, ML, 0.82, 10, 0.4, "CLOUD ALPACAS   ·   SALESFORCE CUSTOMER 360",
         size=11, name=F_SB, color=ORANGE, spacing=2.8)
    text(s, ML, 2.55, 10.6, 2.6, "외부 환경에 흔들리지 않는\n지속 가능한 매출 엔진",
         size=52, name=F_LIGHT, color=WHITE, line_spacing=1.16)
    hline(s, ML, 5.42, 4.6, LINE_DK, 1.0)
    text(s, ML, 5.66, 11, 0.5,
         "Fan 360      Fan Insight      Partner Matching      Sponsorship Sales",
         size=13, name=F_MED, color="AEBECE", spacing=0.6)
    text(s, ML, 6.86, 11, 0.4,
         "Cellsforce · Fan Relationship Management Team        Final Presentation (Draft)",
         size=10, name=F_REG, color="6E8093", spacing=0.4)
    alpaca(s, SW - ML - 2.5, 2.6, 2.5, 2.2, dark=True)
    notes(s, "표지. 한 문장(타이틀) + 서브 1줄 + 팀/날짜. Cloud Alpacas 로고·마스코트 자산을 우측 박스에 배치. "
             "발표자: 오늘 이야기는 '팬을 이해해 매출로 바꾸는 시스템'입니다.")

def s02():
    s = slide()
    kicker(s, "Chapter I · Overview")
    headline(s, "팬은 늘어나는데,\n왜 구단의 매출은\n함께 성장하지 않을까?", y=1.24, size=33, weight=F_LIGHT)
    # qualitative diverging diagram — 방향만, 숫자 없음
    ax_x, ax_y, ax_w = SW - ML - 5.5, 5.7, 4.9
    hline(s, ax_x, ax_y, ax_w, LINE, 0.9)
    text(s, ax_x, ax_y + 0.12, ax_w, 0.3, "시간 →", size=9, name=F_REG, color=FAINT)
    diag(s, ax_x + 0.1, ax_y - 0.15, ax_x + ax_w - 0.1, ax_y - 2.55, NAVY, 1.75)
    diag(s, ax_x + 0.1, ax_y - 1.35, ax_x + ax_w - 0.1, ax_y - 1.05, FAINT, 1.5)
    text(s, ax_x + ax_w - 2.2, ax_y - 3.0, 2.4, 0.3, "팬 수  ↑", size=11, name=F_SB, color=NAVY)
    text(s, ax_x + ax_w - 2.2, ax_y - 0.98, 2.4, 0.3, "구단 매출  ↔ ↓", size=11, name=F_SB, color=MUTED)
    text(s, ML, 6.35, CW, 0.4, "팬 성장이 곧 구단의 지속 가능성은 아니다.",
         size=14, name=F_SB, color=ORANGE_D)
    notes(s, "하나의 큰 질문으로 시작. 두 그래프의 '엇갈림'이 핵심 — 실제 수치가 없으면 방향(↑ / ↔·↓)만. "
             "임의 성장률·매출 수치 만들지 않는다.")
    footer(s, 2)

def s03():
    s = slide()
    kicker(s, "Chapter I · Overview")
    ry = headline(s, "Salesforce 도입 전, 세 가지 문제", size=33)
    rows = [
        ("01", "팬 데이터가 흩어져 있다", "팬의 전체 여정을 하나로 볼 수 없다."),
        ("02", "데이터는 많지만 ACTION이 없다", "팬을 이해해도 다음 행동으로 이어지지 않는다."),
        ("03", "팬덤의 가치를 기업의 기회로 연결할 수 없다", "어떤 기업이 우리 팬덤과 맞는지 판단할 근거가 없다."),
    ]
    y = ry + 0.42
    rh = (6.55 - y) / 3
    for num, head, sub in rows:
        text(s, ML, y + 0.10, 1.2, rh, num, size=30, name=F_LIGHT, color=FAINT)
        text(s, ML + 1.7, y + 0.06, CW - 1.7, 0.5, head, size=19, name=F_MED, color=NAVY)
        text(s, ML + 1.7, y + 0.66, CW - 1.7, 0.5, sub, size=12, name=F_REG, color=MUTED)
        y += rh
        if num != "03":
            hline(s, ML, y, CW, LINE, 0.75)
    notes(s, "Pain Point는 정확히 3개. 큰 숫자 + 강한 headline + 짧은 subline. 긴 설명은 발표자가 말한다.")
    footer(s, 3)

def s04():
    s = slide()
    kicker(s, "Chapter I · Overview")
    headline(s, "반복은 자동화하고,\n사람은 판단과 실행에 집중한다", size=32, weight=F_LIGHT)
    flow(s, ML, 4.15, CW,
         [("Fan Activity", "데이터 발생"), ("Flow", "반복 자동화"),
          ("Insight", "다음 행동 제안"), ("Slack", "담당자 전달"),
          ("Manager Action", "사람이 판단·실행")])
    text(s, ML, 4.95, CW, 0.34,
         "각 단계를  Salesforce · Customer 360 · Flow · Agentforce · Slack  이 잇는다",
         size=9.5, name=F_REG, color=FAINT, align=PP_ALIGN.CENTER)
    text(s, ML, 5.75, CW, 0.5,
         [[("DATA", F_SB, INK), ("      " + ARR + "      ", F_REG, FAINT),
           ("INSIGHT", F_SB, INK), ("      " + ARR + "      ", F_REG, FAINT),
           ("ACTION", F_SB, INK), ("      " + ARR + "      ", F_REG, FAINT),
           ("REVENUE", F_SB, ORANGE)]],
         size=17, align=PP_ALIGN.CENTER, spacing=1.2)
    notes(s, "자동화를 '기술 기능'이 아니라 '업무 방식의 변화'로 설명. 발표자: 사람이 모든 단계를 손으로 잇던 일을 "
             "플랫폼이 잇고, 사람은 판단·실행만 한다.")
    footer(s, 4)

def s05():
    s = slide()
    kicker(s, "Chapter I · Overview")
    headline(s, "B2C에서 시작해 B2B로 이어지는\n하나의 프로젝트", size=32, weight=F_LIGHT)
    flow(s, ML + 1.4, 4.0, CW - 2.8,
         [("B2C", "Fan Relationship"), ("Fan Insight", "연결점 · Bridge"),
          ("B2B", "Sponsorship Sales")], accent_last=False)
    # emphasise the centre pivot
    text(s, ML + 1.4 + (CW - 2.8) / 2 - 1.4, 3.34, 2.8, 0.3, "", size=9)
    text(s, ML, 5.35, CW, 0.34,
         "사라 Fan 360 / Insight    ·    승우 Product · Quote · Campaign    ·    은영 Opportunity    ·    "
         "혜준 Lead    ·    아론 Account · Contact",
         size=10, name=F_MED, color=MUTED, align=PP_ALIGN.CENTER)
    text(s, ML, 6.15, CW, 0.4,
         "각자 자기 구간을 하나의 작은 Salesforce 프로젝트처럼 책임 — 하나의 시나리오로 연결.",
         size=12, name=F_REG, color=FAINT, align=PP_ALIGN.CENTER)
    notes(s, "B2C와 B2B가 별개 프로젝트가 아니라 Fan Insight로 이어진 하나의 흐름임을 보여준다. 팀은 1줄로.")
    footer(s, 5)

def s06():
    s = slide()
    kicker(s, "Chapter II · Demo Scenario", ORANGE)
    ry = headline(s, "From Fan Action to Sponsorship Revenue", size=30)
    text(s, ML, ry + 0.16, CW, 0.32, "지금부터 보실 Demo의 전체 지도 — 10초 안에.",
         size=12, name=F_REG, color=MUTED)
    flow(s, ML, 4.35, CW,
         [("Fan", "QR 참여"), ("Fan 360", "팬 상태"), ("Fan Insight", "기회 발견"),
          ("Partner Matching", "기업 매칭"), ("Lead", "영업 대상"), ("Opportunity", "딜"),
          ("AI Sales", "가속"), ("Closed Won", "매출")],
         accent_last=True, label_size=10)
    text(s, ML, 5.35, CW, 0.4,
         "B2C 3단계   →   Fan Insight 연결점   →   B2B 5단계",
         size=11, name=F_MED, color=MUTED, align=PP_ALIGN.CENTER, spacing=0.8)
    notes(s, "Demo 목차. 'Business Opportunity'라는 말은 쓰지 않는다 — Partner Matching. Fan Insight가 B2C의 "
             "마지막이자 B2B의 출발점.")
    footer(s, 6)

def s07():
    s = slide(NAVY_DEEP)
    text(s, ML, KICK_Y - 0.08, 4, 0.3, "LIVE", size=10, name=F_SB, color=ORANGE, spacing=3.0)
    badge(s, "관객 참여", dark=True)
    text(s, ML, 1.18, 11.5, 1.0, "GAME DAY LIVE", size=46, name=F_LIGHT, color=WHITE, spacing=1.0)
    text(s, ML, 2.28, 11.5, 0.4, "7회말  ·  경기 진행 중", size=14, name=F_MED, color="9DB0C4", spacing=0.8)
    hline(s, ML, 3.05, CW, LINE_DK, 1.0)
    text(s, ML, 3.42, CW - 3.4, 1.4, "오늘의 FAN EVENT\n문태양 선수 퀴즈",
         size=27, name=F_LIGHT, color=WHITE, line_spacing=1.25)
    text(s, ML, 5.02, CW - 3.4, 0.4, "QR을 스캔하고 퀴즈에 참여하세요", size=13, name=F_MED, color="FFCF9E")
    rect(s, SW - ML - 2.0, 3.35, 2.0, 2.0, fill=WHITE, line=None)
    text(s, SW - ML - 2.0, 3.35, 2.0, 2.0, "QR\nPLACEHOLDER", size=10, name=F_SB, color=MUTED,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)
    flow(s, ML, 6.35, CW,
         [("관객 QR 참여",), ("Quiz Entry",), ("Fan Activity",), ("Salesforce",), ("Fan 360",)],
         accent_last=True, sub=False, dark=True, label_size=9.5)
    notes(s, "실제 경기장 전광판처럼. QR destination이 확정 안 됐으면 'QR PLACEHOLDER' 유지 — 가짜 QR을 실제처럼 "
             "표현하지 않는다. Cloud Alpacas 알파카 마스코트를 이벤트 안내 그래픽으로 활용 가능. "
             "발표자: 지금 여러분은 경기장 관객입니다. 전광판에 이벤트가 떴습니다. QR을 찍어보세요. "
             "→ 이 참여가 곧 CRM 데이터가 됩니다. (04_DEMO Scene 1 = 'PPT + 관객 참여')")
    footer(s, 7, dark=True)

def s08():
    demo(8, "B2C", "PPT", "김매니저 · FRM Manager", "우리 팬은 누구인가?",
         "shot", "Fan 360  →  Segment  →  Recommendation Hub 핵심 화면",
         "Fan 360 · Segment · Recommendation Hub        AI 팬 이해 지원        Value  Fan Understanding",
         "발표자: 티켓·굿즈·멤버십·참여가 흩어져 있던 데이터를 한 팬의 프로필·타임라인으로. "
         "실제 Fan 360 / Fan Profile 화면 캡처. (04_DEMO Scene 2 앞부분)")

def s09():
    demo(9, "B2C", "Demo Video", "김매니저 · FRM Manager", "각 팬에게 어떻게 다르게 행동할까?",
         "video", "Target Fan 확인  →  AI 개인화 메시지 생성  →  Review  →  발송   (약 80–90초)",
         "AI Personalized Message        Value  Personalized Fan Engagement",
         "영상 재생 영역이 화면의 중심. 발표자: 팬별 특성을 담당자가 일일이 쓰기 어려웠다 — AI가 생성하고 "
         "담당자가 검토·승인 후 발송(Human-in-the-loop). (04_DEMO 데모 영상 ①)")

def s10():
    s = slide()
    kicker(s, "Chapter II · Demo — Bridge", ORANGE)
    badge(s, "Transition")
    headline(s, "팬의 관심사가\n기업을 찾는 근거가 된다", size=33, weight=F_LIGHT)
    y = 4.05
    text(s, ML, y - 0.42, CW / 2, 0.3, "B2C · 여기까지", size=10, name=F_SB, color=NAVY, spacing=1.8)
    text(s, ML, y, CW / 2, 0.5,
         "Fan Data   " + ARR + "   Fan 360   " + ARR + "   Fan Insight",
         size=13, name=F_MED, color=NAVY)
    text(s, SW / 2 + 0.3, y - 0.42, CW / 2 - 0.3, 0.3, "B2B · 여기부터", size=10, name=F_SB,
         color=ORANGE_D, align=PP_ALIGN.RIGHT, spacing=1.8)
    text(s, SW / 2 + 0.3, y, CW / 2 - 0.3, 0.5,
         "Partner Matching   " + ARR + "   Sponsorship Sales",
         size=13, name=F_MED, color=ORANGE_D, align=PP_ALIGN.RIGHT)
    hline(s, ML, y + 0.85, CW, LINE, 0.75)
    text(s, ML, y + 1.02, CW, 0.6, "FAN INSIGHT", size=30, name=F_LIGHT, color=INK,
         align=PP_ALIGN.CENTER, spacing=3.0)
    oval(s, SW / 2 - 0.06, y + 0.79, 0.12, fill=ORANGE)
    text(s, ML, 6.4, CW, 0.4,
         "20·30대 여성 팬 증가  ·  뷰티 관심 Signal        " + ARR + "        Partner Matching의 근거",
         size=13, name=F_SB, color=ORANGE_D, align=PP_ALIGN.CENTER)
    notes(s, "Demo 전체의 turning point. 여기서 색·발표 톤이 파랑(B2C)→초록(B2B)으로 전환. Fan Insight는 B2B 안이 "
             "아니라 B2C의 마지막 단계. 발표자: 팬의 관심사가 곧 어떤 기업을 찾을지의 근거가 됩니다.")
    footer(s, 10)

def s11():
    demo(11, "B2B", "LIVE", "김매니저 → 이매니저 · FRM → Sponsorship Sales",
         "팬 데이터를 어떻게\nB2B 영업 기회로 연결할까?",
         "live", "Monthly Fan Insight Letter 확인  →  Slack Agent 분석 요청  →  적합한 Sponsorship 방향 탐색",
         "Monthly Fan Insight Letter · Slack Agent        Value  B2C Data → B2B Sales Opportunity",
         "현장 시연. 발표자: B2C에서 쌓인 팬 데이터를 B2B 담당자가 직접 다시 해석해야 했다 — 이제 월간 Fan Insight "
         "Letter와 Slack Agent가 방향을 좁혀준다. 김매니저 → 이매니저 handoff. (04_DEMO Scene 3)")

def s12():
    demo(12, "B2B", "Format TBD", "이매니저 · Sponsorship Sales Manager",
         "이 팬덤과 가장 잘 맞는 기업은?\n— 왜 이 기업인가?",
         "tbd", "Fan Insight  →  Fan Fit  →  Partner Candidate  →  d'Alba  →  Sales Opportunity",
         "Fan Fit · Segment Match · Lead Score        AI Match & Explain        Value  데이터 기반 파트너 발굴",
         "'d'Alba는 팬덤과 높은 적합도를 보여 후보가 됐습니다' — AI는 정답이 아니라 '왜 이 기업인가'를 설명. "
         "Fan Fit / Segment Match(팬덤 적합도) ≠ Lead Score(계약 가능성). 기업 데이터 = OpenDART API. "
         "※ Demo순서 이미지에 S4 열이 없어 재구성 — 포함 여부·표현 방식 팀 확정 필요.",
         media_note="Demo순서 이미지에 S4 미표시 — PPT / VIDEO / LIVE 중 미선택")

def s13():
    demo(13, "B2B", "PPT + 5s Video", "이매니저 / d'Alba · Sponsorship Sales",
         "Sponsor 후보를 어떻게\n실제 Deal로 발전시킬까?",
         "shot5s",
         "Tableau Next Dashboard  →  Lead / Lead Score  →  Account (AI Enrichment)  →  d'Alba OPP",
         "Tableau Next · Lead Score · Account AI Enrichment        Value  Sales Prioritization · Productivity · Data Quality",
         "PPT 위에 약 5초 임베드 영상. 발표자: 유망 Sponsor 판단부터 Account 정보 보완까지 수작업이 많았다 — "
         "우선순위는 Lead Score, 빈 정보는 AI가 DART로 채운다. (04_DEMO Scene 4)")

def s14():
    demo(14, "B2B", "Demo Video", "이매니저 / 김하나 · Sales / d'Alba 담당자",
         "고객은 무엇을 말했는가?",
         "video", "고객 Meeting / Activity  →  기록  →  AI 분석  →  Summary / Signal   (약 3분)",
         "Activity Intelligence        Value  Activity 자산화",
         "영상 재생 영역이 화면의 중심. 발표자: 미팅·대화가 단순 기록으로 남아 다시 읽고 해석해야 했다 — 이제 대화가 "
         "요약·Signal(긍정/위험)로 정리된다. Zoom 연동 · Interaction Intelligence. (04_DEMO Scene 6)")

def s15():
    demo(15, "B2B", "LIVE", "이매니저 / d'Alba · Sponsorship Sales",
         "그래서 무엇을 제안할까?",
         "live",
         "과거 유사 사례 + 현재 OPP + 팬 데이터 + 고객 Activity  →  제안 방향 / Package / Product + 근거",
         "Opportunity Agent        Value  Context 기반 Sales Decision Support",
         "현장 시연. 발표자: 담당자가 컨텍스트를 직접 모아 판단해야 했다 — Agent가 근거와 함께 제안 방향을 제시. "
         "조회·추천은 즉시, 쓰기는 담당자 승인 후. (04_DEMO Scene 7)")

def s16():
    demo(16, "B2B", "PPT", "이매니저 / 김하나 · Sales / d'Alba 담당자",
         "고객의 변화에\n어떻게 대응할까?",
         "shot",
         "새 고객 Activity / 상황  →  AI 선제 분석  →  Negotiation 대응 / 수정안 + 근거  →  Closed Won",
         "Proactive AI · Negotiation Assistant        Value  Proactive Selling",
         "발표자: 고객 반응 변화마다 다시 분석해야 했다 — AI가 변화를 먼저 감지해 수정안을 근거와 함께 제시. "
         "최종 결정·승인은 담당자. ⚠️ 구체 금액은 발표 전 하나로 통일 (04_DEMO 가격 검증). (04_DEMO Scene 8)")

def s17():
    demo(17, "B2B", "Format TBD", "이매니저 / d'Alba · Sponsorship Sales",
         "1년 후, 관계를\n어떻게 다음 매출로 연결할까?",
         "tbd", "'1년 후'  →  단년 계약 종료 임박  →  Partnership Plan 확인  →  d'Alba Upsell",
         "Partnership Plan (논의 필요) · Upsell (논의 필요)        Renewal / Upsell Revenue Expansion (방향)",
         "Demo순서 이미지: 표현 방식·AI 역할·기능 모두 미정. 재계약·장기 Partnership 자동화는 문서상 Future Scope — "
         "구현 확인된 부분만 시연. 발표자: '첫 계약은 매출엔진의 끝이 아니라 시작' — 단, 여기부터는 방향 제시. (04_DEMO Scene 9)",
         media_note="AI 역할·기능·표현 방식 모두 미정 — 팀 확정 필요")

def s18():
    s = slide()
    kicker(s, "Chapter III · Global Best Practices")
    ry = headline(s, "우리가 이렇게 설계한 이유", size=33)
    rows = [
        ("STANDARD FIRST", "표준 Salesforce를 먼저 활용하고, 필요한 부분만 Custom."),
        ("ONE CUSTOMER VIEW", "Fan 360을 중심으로 B2C와 B2B가 같은 데이터를 활용."),
        ("AI + AUTOMATION", "반복은 자동화, AI는 제안, 최종 판단은 사람이 한다."),
    ]
    y = ry + 0.36
    rh = 0.92
    for i, (head, sub) in enumerate(rows):
        text(s, ML, y + 0.04, 0.7, rh, f"0{i+1}", size=15, name=F_LIGHT, color=ORANGE)
        text(s, ML + 1.0, y, 3.6, 0.5, head, size=15, name=F_SB, color=NAVY, spacing=0.8)
        text(s, ML + 4.9, y + 0.02, CW - 4.9, 0.6, sub, size=12, name=F_REG, color=MUTED)
        y += rh
        if i < 2:
            hline(s, ML, y - 0.12, CW, LINE, 0.75)
    flow(s, ML, 6.05, CW,
         [("Fan Activity",), ("Data",), ("Insight",), ("AI / Automation",), ("Human Action",)],
         accent_last=True, sub=False, label_size=10)
    text(s, ML, 6.62, CW, 0.3,
         "17 Custom Objects · 40 Flows · 1 Trigger · 46 LWC · 5 Agentforce Agents · 6 Prompt Templates  (실측)",
         size=8.5, name=F_REG, color=FAINT, align=PP_ALIGN.CENTER)
    notes(s, "4개 기술 나열이 아니라 설계 원칙 3개. 각 원칙은 Headline + 1줄. 구현 수치는 발표자가 필요할 때만 언급. "
             "'왜 이렇게 설계했는가'를 설명하는 페이지.")
    footer(s, 18)

def s19():
    s = slide()
    kicker(s, "Chapter III · Global Best Practices")
    headline(s, "설계 원칙이 실제 업무 흐름으로 작동한다", size=31)
    flow(s, ML, 3.55, CW,
         [("Event / Fan Activity", "데이터 발생"), ("Flow", "반복 자동화"),
          ("Recommendation", "다음 행동 제안"), ("Slack", "담당자 전달"),
          ("Manager Action", "사람이 판단·실행")],
         accent_last=True, label_size=10.5)
    text(s, ML, 4.35, CW, 0.32,
         "Salesforce Platform · Customer 360 · Flow · Agentforce · Slack  이 커넥터처럼 잇는다",
         size=9.5, name=F_REG, color=FAINT, align=PP_ALIGN.CENTER)
    hline(s, ML + 2.2, 5.35, CW - 4.4, LINE, 0.75)
    text(s, ML, 5.6, CW, 1.1,
         [["자동화의 목적은 사람을 대체하는 것이 아니라,"],
          [("담당자가 판단하고 실행할 시간을 ", F_LIGHT, INK), ("확보하는 것.", F_SB, ORANGE_D)]],
         size=19, name=F_LIGHT, color=INK, align=PP_ALIGN.CENTER, line_spacing=1.4)
    notes(s, "기존 Automation 페이지 중심. 단계별 1줄 역할. 긴 설명 금지. 발표자: Flow가 반복을 처리하고 사람은 "
             "판단·실행에 집중한다.")
    footer(s, 19)

def s20():
    s = slide()
    kicker(s, "Chapter IV · Conclusion")
    ry = headline(s, "우리가 만든 것", size=33)
    cols = [
        ("FAN", "팬을 이해하다", "Fan 360 · Personalization", "Fan Lifetime Value  ↑", NAVY),
        ("INSIGHT", "기회를 발견하다", "Fan Insight · Partner Matching", "Personalized Fan Experience  ↑", NAVY),
        ("REVENUE", "매출로 연결하다", "Sponsorship Sales · Opportunity", "Sponsorship Revenue  ↑", ORANGE),
    ]
    y = ry + 0.5
    cw = CW / 3
    for i, (kw, tag, impl, val, col) in enumerate(cols):
        x = ML + i * cw
        if i > 0:
            vline(s, x, y, 3.0, LINE, 0.75)
        px = x + (0.0 if i == 0 else 0.5)
        text(s, px, y, cw - 0.5, 0.8, kw, size=34, name=F_LIGHT, color=col)
        text(s, px, y + 0.92, cw - 0.5, 0.35, tag, size=12, name=F_MED, color=MUTED)
        text(s, px, y + 1.5, cw - 0.6, 0.6, impl, size=11, name=F_REG, color=INK, line_spacing=1.35)
        text(s, px, y + 2.5, cw - 0.6, 0.4, val, size=12.5, name=F_SB, color=ORANGE_D)
    text(s, ML, 6.5, CW, 0.34, "FAN   " + ARR + "   INSIGHT   " + ARR + "   REVENUE   —   하나의 흐름",
         size=11, name=F_MED, color=MUTED, align=PP_ALIGN.CENTER, spacing=1.0)
    notes(s, "What We Built + Business Value 통합. 한 페이지에서 FAN→INSIGHT→REVENUE가 한눈에. 미측정 KPI·ROI 금지.")
    footer(s, 20)

def s21():
    s = slide(NAVY_DEEP)
    rect(s, 0, 0, 0.10, SH, fill=ORANGE)
    text(s, ML, 0.62, 6, 0.3, "NOW · 현재 구현", size=9.5, name=F_SB, color="9DB0C4", spacing=2.0)
    flow(s, ML, 1.12, CW - 0.2,
         [("Fan Data",), ("Insight",), ("Partner Matching",), ("Sponsorship",)],
         accent_last=True, sub=False, dark=True, label_size=9.5)
    text(s, ML, 1.9, 6, 0.3, "FUTURE SCOPE · 미구현", size=9.5, name=F_SB, color="C79A93", spacing=2.0)
    hline(s, ML + 0.1, 2.62, CW - 0.4, "3A3330", 0.9, dash=True)
    for i, lab in enumerate(["Real-time Data", "AI Decision", "Autonomous Action", "Continuous Revenue Growth"]):
        step = (CW - 0.4) / 3
        cx = ML + 0.1 + i * step
        tw = step if 0 < i < 3 else step * 0.9
        text(s, cx - tw / 2, 2.72, tw, 0.4, lab, size=9.5, name=F_REG, color="7E6E6A",
             align=PP_ALIGN.CENTER)
    text(s, ML, 3.7, 10.2, 2.2,
         "팬을 이해하고,\n팬덤의 가치를 발견하고,\n그 가치를 매출로 연결합니다.",
         size=31, name=F_LIGHT, color=WHITE, line_spacing=1.34)
    hline(s, ML, 6.28, 4.6, LINE_DK, 1.0)
    text(s, ML, 6.5, 9, 0.4,
         [[("CLOUD ALPACAS", F_SB, ORANGE), ("      ·      Sustainable Revenue Engine", F_REG, "AEBECE")]],
         size=13, spacing=0.6)
    alpaca(s, SW - ML - 2.4, 3.9, 2.4, 1.9, dark=True)
    notes(s, "발표의 마지막 슬라이드. FUTURE는 현재 구현처럼 보이지 않게(점선·흐리게). 별도 Closing/당첨자 슬라이드 "
             "없음 — 이후 발표자가 '오늘 참여해주신 분들 중…' 하며 LIVE 당첨자 발표 → Q&A로 전환.")
    footer(s, 21, dark=True)

for fn in [s01, s02, s03, s04, s05, s06, s07, s08, s09, s10, s11,
           s12, s13, s14, s15, s16, s17, s18, s19, s20, s21]:
    fn()

prs.save(str(PPTX))
print(f"saved {PPTX.name}  ({len(prs.slides._sldIdLst)} slides)")
